"""利润宝 · 输入数据解析（S4 + S5 扩展）。

支持：
- 从原始 dict 解析为 FinancialData（同义词归并）
- 清洗数字（千分位、括号负数、百分号、货币符号）
- 从 CSV 解析（首列科目，其余列年份）
- 从单个 .xlsx 解析（含"利润表 / 资产负债表 / 科目余额表"三 Sheet）
- 从三个分文件导入并合并
- .xls 给出明确且不误导的错误提示（当前技术栈仅支持 .xlsx）
核心计算仅依赖标准库；Excel 解析按需导入 openpyxl（ADR-004/007）。
"""
from __future__ import annotations

import csv
import os
import re
from typing import Dict, List, Optional, Tuple

from .models import (
    FinancialData, normalize_account_name,
    INCOME_ACCOUNTS, BALANCE_ACCOUNTS, LEDGER_ACCOUNTS,
)

# PDF 预览采集参数（文本型预览：渲染成图会占用数十 MB base64，改为采集文字）
# 2026-08-09：用户要求完整采集（原采样前 20 页会导致分析不完整），改为全部页采集。
_PDF_PREVIEW_MAX_PAGES = 10000     # 安全护栏：不超过 10000 页（审计报告通常 <200 页）
_PDF_TEXT_SAMPLE_PAGES = 5           # 判定文本层覆盖率的采样页数

# 乱码文本层判定（子集化字体/劣质 OCR 层）：
# - 变体字形：带圈数字（⒛22）、罗马数字（Ⅱ）、偏旁（灬丿扌）等
# - 生僻/可疑符号占比过高：乱码层常反复出现 ∽∞冖卜〓㈧⌒ 等字形
_GARBLE_GLYPH_RE = re.compile(
    r"[⒛⒈⒉⒊⒋⒌⒍⒎⒏⒐⒑ⅠⅡⅢⅣⅤⅥⅦⅧⅨⅩ灬丿扌]"
)
_GARBLE_SYMBOL_RE = re.compile(r"[〓㈧㈥∽∞冖卜寸巛廿⌒ˇ≤屮凵ΦΘΝα°^~`]")
_GARBLE_SYMBOL_RATIO = 0.25


def _looks_garbled(text: str) -> bool:
    """文本层是否乱码（子集化字体/劣质 OCR 层）。

    乱码层（如「⒛22年度」「CertiⅡed」或「嘿 剞 羽 斟 释 〓 寮 濠…」）虽含中文，
    但字符错位或由生僻字形反复拼成，直接送大模型会产生错误财务数据；
    识别后应回退 OCR 得到干净文本。
    """
    if _GARBLE_GLYPH_RE.search(text):
        return True
    chars = [ch for ch in text if not ch.isspace()]
    if chars and len(_GARBLE_SYMBOL_RE.findall(text)) / len(chars) >= _GARBLE_SYMBOL_RATIO:
        return True
    # 数字占比过低：正常财报页数字（金额）通常较多
    cjk = len(re.findall(r"[一-鿿]", text))
    digits = len(re.findall(r"[0-9]", text))
    return cjk >= 20 and digits < 3


# ── Sheet 名识别同义词 ────────────────────────────────────────────────────
SHEET_ALIASES = {
    "income": ["利润表", "损益表", "利润及利润分配表", "income", "income statement", "利润"],
    "balance": ["资产负债表", "平衡表", "balance", "balance sheet", "资产"],
    "ledger": ["科目余额表", "余额表", "总账科目余额表", "ledger", "科目余额"],
}


def _match_sheet_kind(name: str) -> Optional[str]:
    """识别 Sheet 名属于哪类报表，返回 income/balance/ledger 或 None。"""
    n = str(name).strip().lower()
    for kind, aliases in SHEET_ALIASES.items():
        for alias in aliases:
            if alias.lower() in n:
                return kind
    return None

# 合规红线：本工具所有建议必须属于合法税务筹划范畴，
# 严禁虚开发票、隐匿收入、虚构成本，严格禁止逃税等违规筹划。
COMPLIANCE_STATEMENT = (
    "本工具仅用于合法税务筹划分析与节税测算，"
    "严禁任何形式的虚开发票、隐匿收入、虚构成本，严格禁止逃税行为。"
)


def clean_number(raw) -> Optional[float]:
    """清洗数字字符串：去千分位 / 空格 / 货币符号 / 百分号；括号与负号表负数。"""
    if raw is None:
        return None
    if isinstance(raw, (int, float)):
        return float(raw)
    s = str(raw).strip()
    if s == "" or s in ("-", "—", "N/A", "无"):
        return None
    neg = False
    if s.startswith("(") and s.endswith(")"):
        neg = True
        s = s[1:-1]
    s = s.replace(",", "").replace("，", "")
    s = re.sub(r"[¥$￥%\s]", "", s)
    if s.startswith("-"):
        neg = True
        s = s[1:]
    if s == "":
        return None
    try:
        val = float(s)
    except ValueError:
        return None
    return -val if neg else val


def _build_statement_table(raw_table: Dict, canonical_accounts: List[str]):
    """将 {科目: {年: 值}} 或 {科目: 值} 归并到规范科目。返回 (table, meta)。"""
    table: Dict[str, Dict[int, Optional[float]]] = {}
    matched = 0
    unmatched: List[str] = []
    for raw_name, val in raw_table.items():
        canon = normalize_account_name(raw_name)
        if canon != raw_name:
            matched += 1
        elif canon not in canonical_accounts:
            unmatched.append(raw_name)
        entry = table.setdefault(canon, {})
        if isinstance(val, dict):
            for yr, v in val.items():
                entry[int(yr)] = clean_number(v)
        else:
            entry["__single__"] = clean_number(val)
    return table, {"matched": matched, "unmatched": unmatched}


def parse_financial_dict(raw: Dict) -> FinancialData:
    """从原始 dict 解析为 FinancialData，完成同义词归并。

    raw 结构：
    {
      "company_name": str, "industry": str,
      "years": [int, ...],
      "income_statement": {科目: {年: 值} | 值},
      "balance_sheet": {...},
      "account_balances": {科目: 值}
    }
    """
    years = [int(y) for y in raw.get("years", [])]
    inc_raw = raw.get("income_statement", {})
    bal_raw = raw.get("balance_sheet", {})
    ledger_raw = raw.get("account_balances", {})

    inc_table, inc_meta = _build_statement_table(inc_raw, INCOME_ACCOUNTS)
    bal_table, bal_meta = _build_statement_table(bal_raw, BALANCE_ACCOUNTS)

    # 单年值落位到年份序列
    for tbl in (inc_table, bal_table):
        for entry in tbl.values():
            if "__single__" in entry and years:
                single = entry.pop("__single__")
                for yr in years:
                    entry.setdefault(yr, single)

    ledger: Dict[str, Optional[float]] = {}
    ledger_unmatched: List[str] = []
    ledger_matched = 0
    for raw_name, val in ledger_raw.items():
        canon = normalize_account_name(raw_name)
        ledger[canon] = clean_number(val)
        if canon != raw_name:
            ledger_matched += 1
        elif canon not in LEDGER_ACCOUNTS:
            ledger_unmatched.append(raw_name)

    # 推断年份（若未提供）
    if not years:
        all_years = set()
        for tbl in (inc_table, bal_table):
            for entry in tbl.values():
                all_years.update(entry.keys())
        years = sorted(all_years)

    meta = {
        "matched": inc_meta["matched"] + bal_meta["matched"] + ledger_matched,
        "unmatched": inc_meta["unmatched"] + bal_meta["unmatched"] + ledger_unmatched,
        "warnings": [],
    }
    if meta["unmatched"]:
        meta["warnings"].append(
            f"未匹配科目 {len(meta['unmatched'])} 项：{meta['unmatched']}"
        )

    return FinancialData(
        company_name=raw.get("company_name", ""),
        industry=raw.get("industry", "制造业"),
        years=years,
        income_statement=inc_table,
        balance_sheet=bal_table,
        account_balances=ledger,
        parsed_meta=meta,
    )


def parse_csv(path: str) -> FinancialData:
    """解析 CSV：首列=科目名，表头其余列=年份。仅支持利润表式单列结构。

    返回 FinancialData（years 来自表头；balance / ledger 为空，供后续扩展）。
    """
    with open(path, "r", encoding="utf-8-sig", newline="") as fh:
        rows = [r for r in csv.reader(fh) if r]
    if not rows:
        raise ValueError("CSV 为空")
    header = rows[0]
    years = [int(h) for h in header[1:] if str(h).strip().isdigit()]
    inc_raw: Dict[str, Dict[int, float]] = {}
    for r in rows[1:]:
        name = r[0].strip()
        if not name:
            continue
        vals = {
            years[i]: clean_number(r[i + 1])
            for i in range(len(years))
            if i + 1 < len(r)
        }
        inc_raw[name] = vals
    return parse_financial_dict({
        "company_name": os.path.basename(path),
        "years": years,
        "income_statement": inc_raw,
    })


# ── Excel 解析（S5 扩展） ─────────────────────────────────────────────────
# 当前技术栈使用 openpyxl，仅支持 .xlsx；.xls（旧版二进制）需要 xlrd 或转
# 换，不在本项目依赖范围内，遇到时给出明确错误提示，不静默失败。

class ParserError(Exception):
    """解析错误基类，供 GUI/上层捕获并展示给用户。"""


def _require_openpyxl():
    try:
        import openpyxl  # noqa: F401
        return openpyxl
    except ImportError as e:
        raise ParserError(
            "未安装 openpyxl，无法解析 Excel。请运行："
            "python3 -m pip install openpyxl"
        ) from e


def _reject_xls(path: str) -> None:
    """对 .xls 文件给出明确且不误导的错误提示。"""
    lower = path.lower()
    if lower.endswith(".xls") and not lower.endswith(".xlsx"):
        raise ParserError(
            f"暂不支持旧版 .xls 文件：{os.path.basename(path)}。"
            "请用 Excel 另存为 .xlsx 后再导入（当前技术栈仅支持 .xlsx）。"
        )


def _read_xlsx_workbook(path: str):
    """读取 .xlsx 工作簿，返回 openpyxl.Workbook（按需加载 openpyxl）。"""
    _reject_xls(path)
    openpyxl = _require_openpyxl()
    try:
        return openpyxl.load_workbook(path, data_only=True, read_only=False)
    except Exception as e:
        raise ParserError(f"无法打开 Excel 文件 {os.path.basename(path)}：{e}") from e


def _detect_year_headers(header_row: List) -> Tuple[List[int], List[int]]:
    """从表头行识别年份列。

    返回 (years, year_col_indices)。表头中连续 4 位数字视为年份。
    """
    years: List[int] = []
    cols: List[int] = []
    for idx, cell in enumerate(header_row):
        s = str(cell).strip() if cell is not None else ""
        if re.fullmatch(r"(19|20)\d{2}", s):
            years.append(int(s))
            cols.append(idx)
    return years, cols


def detect_report_year(path: str, text: str = "") -> Optional[int]:
    """从报告正文或文件名识别报告年度，正文证据优先。"""
    for source in (text, os.path.basename(path)):
        match = re.search(r"((?:19|20)\d{2})\s*年(?:度)?", source or "")
        if match:
            return int(match.group(1))
    return None


def _detect_period_headers(
    header_row: List,
    report_year: Optional[int],
) -> Tuple[List[int], List[int]]:
    """识别明确年份或“本年/上年、期末/期初”等相对期间表头。"""
    years, cols = _detect_year_headers(header_row)
    if years or report_year is None:
        return years, cols

    current_markers = ("本年", "本期", "年末", "期末")
    previous_markers = ("上年", "上期", "年初", "期初")
    for idx, cell in enumerate(header_row):
        value = str(cell).replace(" ", "").strip() if cell is not None else ""
        if any(marker in value for marker in current_markers):
            years.append(report_year)
            cols.append(idx)
        elif any(marker in value for marker in previous_markers):
            years.append(report_year - 1)
            cols.append(idx)
    return years, cols


def _find_header_row(
    ws,
    max_scan: int = 10,
    report_year: Optional[int] = None,
) -> Tuple[int, List[int], List[int]]:
    """在前 max_scan 行中找到含年份的表头行。

    返回 (header_row_idx, years, year_cols)。找不到时返回 (-1, [], [])。
    """
    for r_idx in range(1, min(max_scan + 1, ws.max_row + 1)):
        row = [ws.cell(row=r_idx, column=c).value for c in range(1, ws.max_column + 1)]
        years, cols = _detect_period_headers(row, report_year)
        if years:
            return r_idx, years, cols
    return -1, [], []


def _extract_statement_sheet(
    ws,
    kind: str,
    report_year: Optional[int] = None,
) -> Tuple[Dict[str, Dict[int, Optional[float]]], List[int], List[str]]:
    """从利润表/资产负债表 Sheet 提取 {科目: {年: 值}} 和年份列表。

    返回 (table, years, unmatched_raw_names)。
    """
    header_idx, years, year_cols = _find_header_row(ws, report_year=report_year)
    table: Dict[str, Dict[int, Optional[float]]] = {}
    unmatched: List[str] = []
    if header_idx < 0 or not years:
        return table, [], unmatched
    for r_idx in range(header_idx + 1, ws.max_row + 1):
        name_cell = ws.cell(row=r_idx, column=1).value
        if name_cell is None or str(name_cell).strip() == "":
            continue
        raw_name = str(name_cell).strip()
        canon = normalize_account_name(raw_name)
        canonical_set = INCOME_ACCOUNTS if kind == "income" else BALANCE_ACCOUNTS
        if canon not in canonical_set:
            unmatched.append(raw_name)
            continue
        entry = table.setdefault(canon, {})
        for yi, col_idx in enumerate(year_cols):
            yr = years[yi]
            val = ws.cell(row=r_idx, column=col_idx + 1).value
            entry[yr] = clean_number(val)
    return table, years, unmatched


def _extract_ledger_sheet(ws) -> Tuple[Dict[str, Optional[float]], List[str]]:
    """从科目余额表 Sheet 提取 {规范科目: 值}。

    支持两种列结构：
    1. 首列=科目名称，其余列含"期末余额/期末借方/期末贷方"等表头
    2. 通用：列名含"科目名称"和"期末余额"
    """
    # 找表头行：含"科目名称"或第一行
    name_col = None
    value_col = None
    header_idx = 1
    for r_idx in range(1, min(6, ws.max_row + 1)):
        row = [ws.cell(row=r_idx, column=c).value for c in range(1, ws.max_column + 1)]
        for ci, cell in enumerate(row):
            s = str(cell).strip() if cell is not None else ""
            if s and ("科目" in s and "名称" in s or s.lower() == "account"):
                name_col = ci
            if s and ("期末余额" in s or "期末借方" in s or "余额" == s or s.lower() == "balance"):
                value_col = ci
        if name_col is not None:
            header_idx = r_idx
            break
    if name_col is None:
        # 退化为首列=科目名，第二列=金额
        name_col = 0
        value_col = 1
    if value_col is None:
        value_col = name_col + 1

    ledger: Dict[str, Optional[float]] = {}
    unmatched: List[str] = []
    for r_idx in range(header_idx + 1, ws.max_row + 1):
        name_cell = ws.cell(row=r_idx, column=name_col + 1).value
        if name_cell is None or str(name_cell).strip() == "":
            continue
        raw_name = str(name_cell).strip()
        canon = normalize_account_name(raw_name)
        val = ws.cell(row=r_idx, column=value_col + 1).value
        ledger[canon] = clean_number(val)
        if canon not in LEDGER_ACCOUNTS:
            unmatched.append(raw_name)
    return ledger, unmatched


def parse_excel(path: str, company_name: str = "", industry: str = "制造业") -> FinancialData:
    """解析单个含三 Sheet 的 .xlsx 文件。

    Sheet 名按 SHEET_ALIASES 识别为 income/balance/ledger；
    无法识别的 Sheet 跳过；缺表时记入 parsed_meta["warnings"]。
    """
    wb = _read_xlsx_workbook(path)
    inc_raw: Dict[str, Dict[int, Optional[float]]] = {}
    bal_raw: Dict[str, Dict[int, Optional[float]]] = {}
    ledger_raw: Dict[str, Optional[float]] = {}
    years_all: List[int] = []
    unmatched_all: List[str] = []
    matched_count = 0
    found_kinds = set()
    ignored_sheets: List[str] = []

    for ws in wb.worksheets:
        kind = _match_sheet_kind(ws.title)
        if kind is None:
            ignored_sheets.append(ws.title)
            continue
        if kind == "ledger":
            ledger, unmatched = _extract_ledger_sheet(ws)
            ledger_raw.update(ledger)
            unmatched_all.extend(unmatched)
            found_kinds.add("ledger")
            continue
        table, years, unmatched = _extract_statement_sheet(ws, kind)
        unmatched_all.extend(unmatched)
        if kind == "income":
            inc_raw.update(table)
            found_kinds.add("income")
        else:
            bal_raw.update(table)
            found_kinds.add("balance")
        if years and not years_all:
            years_all = years
        elif years:
            for y in years:
                if y not in years_all:
                    years_all.append(y)
    try:
        wb.close()
    except Exception:
        pass

    warnings: List[str] = []
    for kind, label in [("income", "利润表"), ("balance", "资产负债表"), ("ledger", "科目余额表")]:
        if kind not in found_kinds:
            warnings.append(f"缺表：未在 Excel 中找到{label} Sheet")
    if ignored_sheets:
        warnings.append(f"已忽略未识别的 Sheet：{ignored_sheets}")
    if unmatched_all:
        warnings.append(f"未匹配科目 {len(unmatched_all)} 项：{unmatched_all}")

    raw = {
        "company_name": company_name or os.path.basename(path),
        "industry": industry,
        "years": sorted(years_all),
        "income_statement": inc_raw,
        "balance_sheet": bal_raw,
        "account_balances": ledger_raw,
    }
    data = parse_financial_dict(raw)
    # 合并 warnings（parse_financial_dict 内部也会生成）
    data.parsed_meta["warnings"] = list(data.parsed_meta.get("warnings", [])) + warnings
    data.parsed_meta["excel_path"] = path
    data.parsed_meta["found_kinds"] = sorted(found_kinds)
    return data


# ── Word / PowerPoint 表格解析（Web 版扩展） ─────────────────────────────
# 客户可能以 .docx / .pptx 提供财报表格。将其中表格统一为「网格」后复用
# 上述 openpyxl 提取逻辑，保证与 Excel 解析口径完全一致（Web 对照要求）。

Grid = List[List[object]]


class _GridCell:
    """轻量单元格包装，暴露 .value，使 openpyxl 提取逻辑可直接复用。"""

    __slots__ = ("value",)

    def __init__(self, value: object):
        self.value = value

    def __str__(self) -> str:
        return "" if self.value is None else str(self.value)


class _GridSheet:
    """把 docx/pptx 的表格包装成与 openpyxl worksheet 相近的只读接口。

    仅暴露提取逻辑所需的属性：title / max_row / max_column / cell()。
    cell() 返回 _GridCell，兼容提取函数中的 `.value` 用法。
    """

    def __init__(self, title: str, grid: Grid):
        self.title = title
        self._grid = grid
        self.max_row = len(grid)
        self.max_column = max((len(row) for row in grid), default=0)

    def cell(self, row: int, column: int) -> _GridCell:
        if 0 < row <= self.max_row and 0 < column <= self.max_column:
            r = self._grid[row - 1]
            if 0 < column <= len(r):
                return _GridCell(r[column - 1])
        return _GridCell(None)


def _require_python_docx():
    try:
        import docx  # noqa: F401
        return docx
    except ImportError as e:
        raise ParserError(
            "未安装 python-docx，无法解析 Word 文件。请运行："
            "python3 -m pip install python-docx"
        ) from e


def _require_python_pptx():
    try:
        import pptx  # noqa: F401
        return pptx
    except ImportError as e:
        raise ParserError(
            "未安装 python-pptx，无法解析 PowerPoint 文件。请运行："
            "python3 -m pip install python-pptx"
        ) from e


def _read_docx_tables(path: str) -> List[_GridSheet]:
    """读取 .docx 全部表格为网格列表。"""
    docx = _require_python_docx()
    try:
        doc = docx.Document(path)
    except Exception as e:
        raise ParserError(f"无法打开 Word 文件 {os.path.basename(path)}：{e}") from e
    sheets: List[_GridSheet] = []
    for idx, table in enumerate(doc.tables):
        grid: Grid = []
        for row in table.rows:
            cells = [c.text for c in row.cells]
            grid.append(cells)
        sheets.append(_GridSheet(f"doc_table_{idx}", grid))
    if not sheets:
        raise ParserError(f"Word 文件 {os.path.basename(path)} 中未找到表格")
    return sheets


def _read_pptx_tables(path: str) -> List[_GridSheet]:
    """读取 .pptx 全部幻灯片内表格为网格列表。"""
    pptx = _require_python_pptx()
    try:
        prs = pptx.Presentation(path)
    except Exception as e:
        raise ParserError(f"无法打开 PowerPoint 文件 {os.path.basename(path)}：{e}") from e
    sheets: List[_GridSheet] = []
    idx = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            grid: Grid = []
            for row in shape.table.rows:
                grid.append([c.text for c in row.cells])
            sheets.append(_GridSheet(f"ppt_table_{idx}", grid))
            idx += 1
    if not sheets:
        raise ParserError(f"PowerPoint 文件 {os.path.basename(path)} 中未找到表格")
    return sheets


# OCR 采集参数（RapidOCR 惰性加载，未安装时静默降级）
# 2026-08-09：用户要求完整采集，不再限制 OCR 页数上限（扫描件全量识别）
_OCR_MAX_PAGES = 10000          # 安全护栏：全量 OCR（不再只识别前 12 页）
_ocr_engine_holder: list = []  # 惰性单例容器（元素 0 为引擎或 None）


def _get_ocr_engine():
    """惰性加载 RapidOCR 引擎；不可用时返回 None（不阻塞预览）。"""
    if _ocr_engine_holder:
        return _ocr_engine_holder[0]
    try:
        from rapidocr_onnxruntime import RapidOCR
        engine = RapidOCR()
    except Exception:
        engine = None
    _ocr_engine_holder.append(engine)
    return engine


def _ocr_pil_text(engine, pil_image) -> str:
    """对 PIL 图片 OCR，返回按行拼接的文本；失败返回空串。"""
    try:
        import numpy as np
        arr = np.array(pil_image.convert("RGB"))
        result, _ = engine(arr)
        if not result:
            return ""
        lines = [str(item[1]).strip() for item in result if item and len(item) > 1 and item[1]]
        return "\n".join(lines)
    except Exception:
        return ""


def _render_page_pil(pdf, page_index: int, scale: float = 1.5):
    """渲染 PDF 单页为 PIL 图片（供 OCR/预览采样）。失败返回 None。"""
    try:
        import pypdfium2 as pdfium
    except ImportError:
        return None
    try:
        return pdf[page_index].render(scale=scale).to_pil()
    except Exception:
        return None


def _page_text_layer(path: str, page_index: int) -> str:
    """取 PDF 单页文本层（pdfplumber）。失败/为空返回空串。"""
    pdfplumber = _require_pdfplumber()
    try:
        with pdfplumber.open(path) as pdf:
            if page_index >= len(pdf.pages):
                return ""
            return (pdf.pages[page_index].extract_text() or "").strip()
    except Exception:
        return ""


def _pdf_text_coverage(path: str, sample_pages: int) -> float:
    """按前 sample_pages 页文本层有效字符数评估覆盖率（0.0~1.0）。

    用于区分扫描件（无文本层）与普通 PDF：只计算非空白、非版式噪声字符。
    """
    total = 0
    filled = 0
    for pi in range(sample_pages):
        txt = _page_text_layer(path, pi)
        if not txt:
            continue
        total += 1
        # 中文字符 + 数字 + 拉丁字母视为有效内容；版式符号/标点不计
        valid = len(re.findall(r"[一-鿿0-9A-Za-z]", txt))
        if valid >= 20:
            filled += 1
    if total == 0:
        return 0.0
    return filled / sample_pages


def _classify_pdf_type(path: str) -> str:
    """按文本层覆盖率判定 PDF 类型：scan / text / mixed。

    scan：绝大多数页无文本层（扫描件，需 AI 解析）；
    text：前若干页均有足量文本层；
    mixed：部分有、部分无。
    """
    coverage = _pdf_text_coverage(path, _PDF_TEXT_SAMPLE_PAGES)
    if coverage == 0.0:
        return "scan"
    if coverage >= 0.8:
        return "text"
    return "mixed"


def _require_pdfplumber():
    try:
        import pdfplumber  # noqa: F401
        return pdfplumber
    except ImportError as e:
        raise ParserError(
            "未安装 pdfplumber，无法解析 PDF 文件。请运行："
            "python3 -m pip install pdfplumber"
        ) from e


def _rebuild_pdf_grid(page, row_tol: float = 6.0, col_gap: float = 14.0) -> Grid:
    """用 words 坐标重建表格网格（不依赖网格线）。

    很多报表 PDF（金蝶/用友导出）的表格线是背景填充而非矢量线，
    extract_tables() 识别失败。改用词坐标按行分组、按 x 间隙分列，
    能稳定重建 科目名|本期|上期 这类报表结构。

    row_tol：同一行词允许的 top 最大间隙（像素）。PDF 中同一行词因
    字体基线对齐差异 top 可差 3-5px，需用贪心聚类合并而非固定桶。
    """
    words = page.extract_words()
    if not words:
        return []
    # 贪心聚类：按 top 排序，间隙 < row_tol 的词并入同一行
    words = sorted(words, key=lambda w: w["top"])
    rows: List[List] = []
    cur_row = [words[0]]
    for w in words[1:]:
        if w["top"] - cur_row[-1]["top"] <= row_tol:
            cur_row.append(w)
        else:
            rows.append(cur_row)
            cur_row = [w]
    rows.append(cur_row)

    grid: Grid = []
    for line in rows:
        line = sorted(line, key=lambda w: w["x0"])
        cells: List[str] = []
        cur = [line[0]]
        for w in line[1:]:
            if w["x0"] - cur[-1]["x1"] > col_gap:
                cells.append(" ".join(x["text"] for x in cur))
                cur = [w]
            else:
                cur.append(w)
        cells.append(" ".join(x["text"] for x in cur))
        grid.append(cells)
    return grid


def _read_pdf_tables(path: str) -> List[_GridSheet]:
    """读取 .pdf 全部页面中可识别的表格为网格列表。

    优先使用 pdfplumber 的 extract_tables()（有网格线的 PDF）；
    识别为空时回退到坐标重建（报表线为背景填充的 PDF）。
    纯文字扫描版 PDF（无文本层）无法提取，预览中会给出提示。
    """
    pdfplumber = _require_pdfplumber()
    try:
        pdf = pdfplumber.open(path)
    except Exception as e:
        raise ParserError(f"无法打开 PDF 文件 {os.path.basename(path)}：{e}") from e
    sheets: List[_GridSheet] = []
    try:
        for pi, page in enumerate(pdf.pages):
            try:
                tables = page.extract_tables()
            except Exception:
                tables = []
            if not tables:
                rebuilt = _rebuild_pdf_grid(page)
                if rebuilt:
                    tables = [rebuilt]
            for ti, grid in enumerate(tables):
                if not grid:
                    continue
                # 过滤完全空行，避免噪声
                cleaned: Grid = [row for row in grid if any((c or "").strip() for c in row)]
                if not cleaned:
                    continue
                sheets.append(_GridSheet(f"pdf_page{pi + 1}_table{ti + 1}", cleaned))
    finally:
        pdf.close()
    if not sheets:
        fname = os.path.basename(path)
        if _pdf_text_coverage(path, _PDF_TEXT_SAMPLE_PAGES) == 0.0:
            raise ParserError(
                f"PDF 文件 {fname} 为扫描件（无文本层），无法直接提取表格。"
                "请先在预览区「配置 AI」后重试，或改用 Excel/CSV/Word 导入。"
            )
        raise ParserError(
            f"PDF 文件 {fname} 中未找到可解析的表格。"
            "若为扫描件或乱码文本（子集化字体），请配置 AI 后重试，"
            "或改用 Excel/CSV/Word 导入。"
        )
    return sheets


def _parse_grid_sheets(sheets: List[_GridSheet], company_name: str, industry: str,
                       source_label: str) -> FinancialData:
    """从网格列表解析统一 FinancialData（复用 openpyxl 提取逻辑）。"""
    inc_raw: Dict[str, Dict[int, Optional[float]]] = {}
    bal_raw: Dict[str, Dict[int, Optional[float]]] = {}
    ledger_raw: Dict[str, Optional[float]] = {}
    years_all: List[int] = []
    unmatched_all: List[str] = []
    found_kinds = set()
    ignored: List[str] = []
    report_year = detect_report_year(source_label)

    for gs in sheets:
        kind = _match_sheet_kind(gs.title) or _sniff_sheet_kind(gs)
        if kind is None:
            ignored.append(gs.title)
            continue
        if kind == "ledger":
            ledger, unmatched = _extract_ledger_sheet(gs)
            ledger_raw.update(ledger)
            unmatched_all.extend(unmatched)
            found_kinds.add("ledger")
            continue
        table, years, unmatched = _extract_statement_sheet(
            gs, kind, report_year=report_year,
        )
        unmatched_all.extend(unmatched)
        if kind == "income":
            inc_raw.update(table)
            found_kinds.add("income")
        else:
            bal_raw.update(table)
            found_kinds.add("balance")
        if years and not years_all:
            years_all = years
        elif years:
            for y in years:
                if y not in years_all:
                    years_all.append(y)

    warnings: List[str] = []
    for kind, label in [("income", "利润表"), ("balance", "资产负债表"), ("ledger", "科目余额表")]:
        if kind not in found_kinds:
            warnings.append(f"缺表：未在{source_label}中找到{label}表格")
    if ignored:
        warnings.append(f"已忽略无法识别类型的表格：{ignored}")
    if unmatched_all:
        warnings.append(f"未匹配科目 {len(unmatched_all)} 项：{unmatched_all}")

    raw = {
        "company_name": company_name or os.path.basename(source_label),
        "industry": industry,
        "years": sorted(years_all),
        "income_statement": inc_raw,
        "balance_sheet": bal_raw,
        "account_balances": ledger_raw,
    }
    data = parse_financial_dict(raw)
    data.parsed_meta["warnings"] = list(data.parsed_meta.get("warnings", [])) + warnings
    data.parsed_meta["source_path"] = source_label
    data.parsed_meta["found_kinds"] = sorted(found_kinds)
    return data


def _sniff_sheet_kind(gs: _GridSheet) -> str:
    """无标题网格按内容特征推断类型（docx/pptx 表格没有 Sheet 名）。

    判定顺序：
    1. 表头/首行含「科目名称」「期末余额」→ ledger
    2. 含「利润表」「资产负债表」等字样 → income/balance
    3. 首行含年份表头时，按第一列科目归并到哪类规范科目判定
    """
    # 首行特征：科目余额表
    header = " ".join(str(gs.cell(1, c) or "") for c in range(1, min(gs.max_column + 1, 6)))
    if "科目名称" in header or "期末余额" in header:
        return "ledger"

    # 内容含表名关键字
    for r_idx in range(1, min(gs.max_row + 1, 6)):
        row_text = " ".join(str(gs.cell(r_idx, c) or "") for c in range(1, gs.max_column + 1))
        if "利润表" in row_text or "损益表" in row_text:
            return "income"
        if "资产负债表" in row_text:
            return "balance"

    # 按科目归并判定：扫描科目列，看命中利润表还是资产负债表科目更多
    # （不依赖年份表头——「本期/上期」等报表列同样适用）
    income_hits = balance_hits = 0
    for r_idx in range(2, min(gs.max_row + 1, 40)):
        name = str(gs.cell(r_idx, 1) or "").strip()
        if not name:
            continue
        canon = normalize_account_name(name)
        if canon in INCOME_ACCOUNTS:
            income_hits += 1
        elif canon in BALANCE_ACCOUNTS:
            balance_hits += 1
    if income_hits or balance_hits:
        if income_hits >= balance_hits:
            return "income"
        return "balance"
    return ""


def parse_docx(path: str, company_name: str = "", industry: str = "制造业") -> FinancialData:
    """解析 .docx 中的财报表格为统一 FinancialData。"""
    sheets = _read_docx_tables(path)
    return _parse_grid_sheets(sheets, company_name, industry, os.path.basename(path))


def parse_pptx(path: str, company_name: str = "", industry: str = "制造业") -> FinancialData:
    """解析 .pptx 中的财报表格为统一 FinancialData。"""
    sheets = _read_pptx_tables(path)
    return _parse_grid_sheets(sheets, company_name, industry, os.path.basename(path))


def parse_pdf(path: str, company_name: str = "", industry: str = "制造业") -> FinancialData:
    """解析 .pdf 中的财报表格为统一 FinancialData。"""
    sheets = _read_pdf_tables(path)
    return _parse_grid_sheets(sheets, company_name, industry, os.path.basename(path))


def parse_excel_files(
    income_path: Optional[str] = None,
    balance_path: Optional[str] = None,
    ledger_path: Optional[str] = None,
    company_name: str = "",
    industry: str = "制造业",
) -> FinancialData:
    """解析三个分文件并合并为统一 FinancialData。

    每个文件可以是单 Sheet 的 .xlsx，也可以是 CSV（首列科目，其余列年份）。
    至少需要 income_path；其余可缺。
    """
    if not income_path:
        raise ParserError("至少需要提供利润表文件（income_path）")

    def _merge(target: Dict, source: Dict) -> None:
        for k, v in source.items():
            if k not in target:
                target[k] = v
            elif isinstance(v, dict) and isinstance(target[k], dict):
                target[k].update(v)

    years_all: List[int] = []
    inc_raw: Dict[str, Dict[int, Optional[float]]] = {}
    bal_raw: Dict[str, Dict[int, Optional[float]]] = {}
    ledger_raw: Dict[str, Optional[float]] = {}
    warnings: List[str] = []

    def _parse_one(path: str, kind: str) -> None:
        nonlocal years_all
        if not os.path.exists(path):
            raise ParserError(f"文件不存在：{path}")
        lower = path.lower()
        if lower.endswith(".csv"):
            fd = parse_csv(path)
            if kind == "income":
                _merge(inc_raw, fd.income_statement)
            elif kind == "balance":
                _merge(bal_raw, fd.balance_sheet)
            for y in fd.years:
                if y not in years_all:
                    years_all.append(y)
            return
        wb = _read_xlsx_workbook(path)
        for ws in wb.worksheets:
            if kind == "ledger":
                ledger, _ = _extract_ledger_sheet(ws)
                ledger_raw.update(ledger)
            else:
                table, years, _ = _extract_statement_sheet(ws, kind)
                if kind == "income":
                    _merge(inc_raw, table)
                else:
                    _merge(bal_raw, table)
                for y in years:
                    if y not in years_all:
                        years_all.append(y)
        try:
            wb.close()
        except Exception:
            pass

    _parse_one(income_path, "income")
    if balance_path:
        _parse_one(balance_path, "balance")
    else:
        warnings.append("缺表：未提供资产负债表文件")
    if ledger_path:
        _parse_one(ledger_path, "ledger")
    else:
        warnings.append("缺表：未提供科目余额表文件")

    raw = {
        "company_name": company_name or os.path.basename(income_path),
        "industry": industry,
        "years": sorted(years_all),
        "income_statement": inc_raw,
        "balance_sheet": bal_raw,
        "account_balances": ledger_raw,
    }
    data = parse_financial_dict(raw)
    data.parsed_meta["warnings"] = list(data.parsed_meta.get("warnings", [])) + warnings
    return data


def _guess_file_year(name: str) -> Optional[int]:
    """从文件名识别审计报告年度，如「2022年审计报告.pdf」→ 2022。

    识别四位数年份后跟「年」的模式；取不到返回 None。
    """
    m = re.search(r"(19|20)\d{2}\s*年", name)
    if m:
        return int(re.search(r"(19|20)\d{2}", m.group(0)).group(0))
    return None


def merge_years(*datas: FinancialData) -> FinancialData:
    """合并多份 FinancialData（各是一整年的完整审计报告）为三年数据集。

    按 科目+年份 合并 income_statement / balance_sheet / account_balances，
    years 取并集。company_name / industry 取首个非空。
    科目余额表（account_balances）带年份维度，各年各自保留。
    """
    if not datas:
        raise ParserError("合并数据为空")
    merged = FinancialData(
        company_name=next((d.company_name for d in datas if d.company_name), ""),
        industry=next((d.industry for d in datas if d.industry), "制造业"),
    )

    def _merge_year_dict(target: Dict[str, Dict[int, Optional[float]]],
                         source: Dict[str, Dict[int, Optional[float]]]) -> None:
        for acc, year_vals in source.items():
            if not year_vals:
                continue
            target.setdefault(acc, {})
            for yr, val in year_vals.items():
                if yr in ("__single__",):
                    continue
                target[acc][int(yr)] = val

    years_all: List[int] = []
    warnings: List[str] = []
    for d in datas:
        _merge_year_dict(merged.income_statement, d.income_statement)
        _merge_year_dict(merged.balance_sheet, d.balance_sheet)
        _merge_year_dict(merged.account_balances, d.account_balances)
        for y in d.years:
            if y not in years_all:
                years_all.append(y)
        meta = d.parsed_meta or {}
        warnings.extend(meta.get("warnings", []))

    merged.years = sorted(years_all)
    merged.parsed_meta = {
        "matched": sum(len(d.income_statement) for d in datas),
        "unmatched": [],
        "warnings": warnings,
        "merged_files": len(datas),
    }
    return merged


def make_empty_data(company_name: str = "", industry: str = "制造业") -> FinancialData:
    """构造空占位 FinancialData（用于纯扫描件 PDF 等无法结构化解析的场景）。

    保留企业/行业信息，指标与科目为空；OCR 文本仍可被预览采集供 AI 使用。
    """
    return FinancialData(
        company_name=company_name,
        industry=industry or "制造业",
        years=[],
        parsed_meta={"warnings": ["扫描件无法结构化解析，已用 OCR 文本供 AI 使用"]},
    )


def parse_smart(
    path: Optional[str] = None,
    income_path: Optional[str] = None,
    balance_path: Optional[str] = None,
    ledger_path: Optional[str] = None,
    company_name: str = "",
    industry: str = "制造业",
) -> FinancialData:
    """智能解析：单文件优先，其次分文件，CSV/Excel/Word/PPT/PDF 自动识别。"""
    if path:
        lower = path.lower()
        if lower.endswith(".csv"):
            return parse_csv(path)
        if lower.endswith(".xls") and not lower.endswith(".xlsx"):
            _reject_xls(path)
        if lower.endswith(".docx"):
            return parse_docx(path, company_name=company_name, industry=industry)
        if lower.endswith(".pptx"):
            return parse_pptx(path, company_name=company_name, industry=industry)
        if lower.endswith(".pdf"):
            return parse_pdf(path, company_name=company_name, industry=industry)
        return parse_excel(path, company_name=company_name, industry=industry)
    if income_path or balance_path or ledger_path:
        return parse_excel_files(
            income_path=income_path,
            balance_path=balance_path,
            ledger_path=ledger_path,
            company_name=company_name,
            industry=industry,
        )
    raise ParserError("未提供任何输入文件路径")


# ── 文件内容预览提取（Web 版导入页展示） ─────────────────────────────────
# 返回 { "name", "kind", "sections": [{ "title", "grid": [[..]] }] }，
# grid 元素统一为字符串，供前端原样渲染。

def preview_file(path: str) -> dict:
    """提取文件可预览内容（表格 + 非表格文本）。支持 xlsx/csv/docx/pptx/pdf。

    返回:
        {
          "name": 文件名,
          "sections": [ { "title": 节标题, "grid": [[str, ...], ...] } ],
          "notes": [ ...非表格文本片段... ],
        }
    """
    lower = path.lower()
    name = os.path.basename(path)

    if lower.endswith(".csv"):
        return _preview_csv(path, name)
    if lower.endswith(".xlsx"):
        return _preview_xlsx(path, name)
    if lower.endswith(".docx"):
        return _preview_docx(path, name)
    if lower.endswith(".pptx"):
        return _preview_pptx(path, name)
    if lower.endswith(".pdf"):
        return _preview_pdf(path, name)
    if lower.endswith(".xls") and not lower.endswith(".xlsx"):
        _reject_xls(path)
    raise ParserError(f"暂不支持预览：{name}")


def _cells_to_str(grid: Grid) -> List[List[str]]:
    return [[("" if c is None else str(c).strip()) for c in row] for row in grid]


def _preview_csv(path: str, name: str) -> dict:
    with open(path, "r", encoding="utf-8-sig", newline="") as fh:
        rows = [r for r in csv.reader(fh) if r]
    return {"name": name, "kind": "csv", "sections": [{"title": "CSV 内容", "grid": _cells_to_str(rows)}], "notes": []}


def _preview_xlsx(path: str, name: str) -> dict:
    wb = _read_xlsx_workbook(path)
    sections = []
    for ws in wb.worksheets:
        grid = []
        for row in ws.iter_rows(values_only=True):
            grid.append(list(row))
        sections.append({"title": ws.title, "grid": _cells_to_str(grid)})
    try:
        wb.close()
    except Exception:
        pass
    return {"name": name, "kind": "xlsx", "sections": sections, "notes": []}


def _preview_docx(path: str, name: str) -> dict:
    docx = _require_python_docx()
    try:
        doc = docx.Document(path)
    except Exception as e:
        raise ParserError(f"无法打开 Word 文件 {name}：{e}") from e
    sections = []
    for idx, table in enumerate(doc.tables):
        grid = [[c.text for c in row.cells] for row in table.rows]
        sections.append({"title": f"表格 {idx + 1}", "grid": _cells_to_str(grid)})
    notes = [p.text for p in doc.paragraphs if p.text.strip()]
    return {"name": name, "kind": "docx", "sections": sections, "notes": notes[:50]}


def _preview_pptx(path: str, name: str) -> dict:
    pptx = _require_python_pptx()
    try:
        prs = pptx.Presentation(path)
    except Exception as e:
        raise ParserError(f"无法打开 PowerPoint 文件 {name}：{e}") from e
    sections = []
    notes = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                grid = [[c.text for c in row.cells] for row in shape.table.rows]
                sections.append({"title": f"幻灯片 {si + 1} 表格", "grid": _cells_to_str(grid)})
            elif shape.has_text_frame and shape.text_frame.text.strip():
                notes.append(f"[幻灯片 {si + 1}] {shape.text_frame.text.strip()}")
    return {"name": name, "kind": "pptx", "sections": sections, "notes": notes[:50]}


def _preview_pdf(path: str, name: str) -> dict:
    """预览 PDF：按页采集文本层与 OCR 文字，返回文本型预览（不再渲染图片）。

    返回结构：
        {
          "name", "kind": "pdf", "pdf_type": "scan|text|mixed",
          "images": [],           # 文本型预览，不再携带 base64 图片
          "sections": [],         # 兼容字段，PDF 不用表格预览
          "notes": [逐页文本/OCR 片段...],
        }
    - 渲染成 base64 大图会占用数十 MB 响应并拖垮浏览器，已改为采集文字
    - 文本层页加入 [第 N 页] 前缀，OCR 页加入 [第 N 页 OCR] 前缀（供 AI 使用）
    - pdf_type 用于前端提示：扫描件需配置 AI 才能导入
    """
    try:
        import pypdfium2 as pdfium
    except ImportError as e:
        raise ParserError(
            "未安装 pypdfium2，无法预览 PDF 文件。请运行："
            "python3 -m pip install pypdfium2"
        ) from e

    try:
        pdf = pdfium.PdfDocument(path)
    except Exception as e:
        raise ParserError(f"无法打开 PDF 文件 {name}：{e}") from e

    notes: list[str] = []
    ocr = _get_ocr_engine()
    total = len(pdf)
    max_pages = min(total, _PDF_PREVIEW_MAX_PAGES)
    ocr_pages = min(total, _OCR_MAX_PAGES)
    try:
        for pi in range(max_pages):
            # 1) 文本层
            layer = _page_text_layer(path, pi)
            if layer:
                # 截断超长文本层，避免 notes 过大（4000 字符≈整页正文）
                layer = layer[:4000]
                if _looks_garbled(layer):
                    # 文本层乱码（子集化字体等）：不能直接喂 AI，回退 OCR
                    ocr_text = ""
                    if ocr is not None and pi < ocr_pages:
                        pil = _render_page_pil(pdf, pi)
                        if pil is not None:
                            ocr_text = _ocr_pil_text(ocr, pil) or ""
                    if ocr_text:
                        notes.append(f"[第 {pi + 1} 页 OCR·文本层乱码]\n{ocr_text[:4000]}")
                    else:
                        notes.append(
                            f"[第 {pi + 1} 页 文本层乱码]\n"
                            "（本页文本层为乱码且未采集到可用 OCR 文字）"
                        )
                    continue
                notes.append(f"[第 {pi + 1} 页]\n{layer}")
                continue  # 文本层页不再重复 OCR，节省耗时
            # 2) 无文本层：OCR 采集（供 AI 整理/扫描件解析）
            if ocr is None or pi >= ocr_pages:
                continue
            pil = _render_page_pil(pdf, pi)
            if pil is None:
                continue
            ocr_text = _ocr_pil_text(ocr, pil)
            if ocr_text:
                notes.append(f"[第 {pi + 1} 页 OCR]\n{ocr_text[:4000]}")
    finally:
        pdf.close()

    if total > max_pages:
        notes.insert(0, f"共 {total} 页，预览采样前 {max_pages} 页文字")
    else:
        notes.insert(0, f"共 {total} 页，已完整采集全部页文字")
    return {
        "name": name,
        "kind": "pdf",
        "pdf_type": _classify_pdf_type(path),
        "sections": [],
        "images": [],
        "notes": notes,
    }
