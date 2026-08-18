"""利润宝 · Web 导入扩展测试（docx/pptx/pdf 解析与 5 格式预览）。

覆盖：Word/PPT/PDF 表格解析、多格式导入闭环、预览端点内容、错误降级。
全部使用 TestClient 与临时生成文件，不依赖浏览器。
"""

from __future__ import annotations

import io
import importlib

import pytest
from fastapi.testclient import TestClient

app_mod = importlib.import_module("web_backend.CO_app_WB-CO-TR-20260805160732")
import_mod = importlib.import_module("web_backend.CO_import_WB-CO-TR-20260805160732")
session_mod = importlib.import_module("web_backend.CO_session_WB-CO-TR-20260805160732")

CLIENT = TestClient(app_mod.create_app())


@pytest.fixture(autouse=True)
def _reset_session():
    session_mod.clear()
    yield
    session_mod.clear()


# ── 构造 docx / pptx 测试文件 ───────────────────────────────────────────

def _build_docx_bytes() -> bytes:
    """生成含三张财报表格的 Word 文件。"""
    import docx

    from data.make_sample import INCOME, BALANCE, LEDGER

    doc = docx.Document()

    def _add_table(rows):
        t = doc.add_table(rows=len(rows), cols=len(rows[0]))
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                t.cell(ri, ci).text = str(val)

    _add_table([["项目", "2021", "2022", "2023"]]
               + [[acc] + [INCOME[acc][y] for y in (2021, 2022, 2023)] for acc in INCOME])
    _add_table([["项目", "2021", "2022", "2023"]]
               + [[acc] + [BALANCE[acc][y] for y in (2021, 2022, 2023)] for acc in BALANCE])
    _add_table([["科目名称", "期末余额"]] + [[acc, val] for acc, val in LEDGER.items()])

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def _build_pptx_bytes() -> bytes:
    """生成含三张财报表格的 PowerPoint 文件。"""
    from pptx import Presentation

    from data.make_sample import INCOME, BALANCE, LEDGER

    prs = Presentation()

    def _add_table(slide, rows):
        shape = slide.shapes.add_table(len(rows), len(rows[0]), 10, 10, 100, 100)
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                shape.table.cell(ri, ci).text = str(val)

    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    _add_table(s1, [["项目", "2021", "2022", "2023"]]
               + [[acc] + [INCOME[acc][y] for y in (2021, 2022, 2023)] for acc in INCOME])
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    _add_table(s2, [["项目", "2021", "2022", "2023"]]
               + [[acc] + [BALANCE[acc][y] for y in (2021, 2022, 2023)] for acc in BALANCE])
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    _add_table(s3, [["科目名称", "期末余额"]] + [[acc, val] for acc, val in LEDGER.items()])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _xlsx_bytes() -> bytes:
    from data.make_sample import write_sample_xlsx
    import os

    path = write_sample_xlsx()
    try:
        with open(path, "rb") as fh:
            return fh.read()
    finally:
        os.remove(path)


def _build_pdf_bytes() -> bytes:
    """生成嵌入中文字体的财报表格 PDF（matplotlib 渲染，保证 ToUnicode 可提取）。

    使用 PdfPages 写多页：利润表 / 资产负债表 / 科目余额表。
    """
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.backends.backend_pdf import PdfPages

    from data.make_sample import INCOME, BALANCE, LEDGER

    plt.rcParams["font.sans-serif"] = ["Heiti SC", "STHeiti", "Hiragino Sans GB"]
    plt.rcParams["axes.unicode_minus"] = False

    def _draw_table(fig, rows):
        ax = fig.add_subplot(111)
        ax.axis("off")
        t = ax.table(cellText=rows, loc="center")
        t.auto_set_font_size(False)
        t.set_fontsize(9)

    buf = io.BytesIO()
    with PdfPages(buf) as pdf:
        fig = plt.figure(figsize=(9, 7))
        _draw_table(fig, [["项目", "2021", "2022", "2023"]]
                    + [[a] + [INCOME[a][y] for y in (2021, 2022, 2023)] for a in INCOME])
        pdf.savefig(fig)
        plt.close(fig)

        fig = plt.figure(figsize=(9, 7))
        _draw_table(fig, [["项目", "2021", "2022", "2023"]]
                    + [[a] + [BALANCE[a][y] for y in (2021, 2022, 2023)] for a in BALANCE])
        pdf.savefig(fig)
        plt.close(fig)

        fig = plt.figure(figsize=(6, 5))
        _draw_table(fig, [["科目名称", "期末余额"]] + [[a, v] for a, v in LEDGER.items()])
        pdf.savefig(fig)
        plt.close(fig)
    buf.seek(0)
    return buf.read()


# ── docx / pptx 解析 ───────────────────────────────────────────────────

def test_parse_docx_tables():
    import docx
    from core.parser import parse_docx

    doc = docx.Document(io.BytesIO(_build_docx_bytes()))
    # 三张表格
    assert len(doc.tables) == 3
    # 第一张表格含年份表头
    assert doc.tables[0].cell(0, 1).text == "2021"


def test_import_docx():
    r = CLIENT.post("/api/import", files={"files": ("财报.docx", io.BytesIO(_build_docx_bytes()), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}, data={"company_name": "文档厂"})
    assert r.status_code == 200, r.text
    body = r.json()
    assert body["summary"]["company_name"] == "文档厂"
    assert body["summary"]["years"] == [2021, 2022, 2023]
    # docx 解析后会计量
    data = session_mod.get_data()
    assert data is not None
    assert data.income_statement.get("营业收入", {}).get(2023) == 15_200_000
    assert data.account_balances.get("咨询服务费") == {2023: 980_000.0}
    # 预览含表格
    assert body["previews"], "导入响应应含预览"
    assert body["previews"][0]["kind"] == "docx"
    assert body["previews"][0]["sections"][0]["grid"][0][0] == "项目"


def test_import_pptx():
    r = CLIENT.post("/api/import", files={"files": ("财报.pptx", io.BytesIO(_build_pptx_bytes()), "application/vnd.openxmlformats-officedocument.presentationml.presentation")}, data={"company_name": "幻灯厂"})
    assert r.status_code == 200, r.text
    body = r.json()
    assert body["summary"]["years"] == [2021, 2022, 2023]
    data = session_mod.get_data()
    assert data is not None
    assert data.income_statement.get("营业收入", {}).get(2023) == 15_200_000
    assert body["previews"][0]["kind"] == "pptx"
    assert body["previews"][0]["sections"][0]["grid"][0][0] == "项目"


# ── 多格式导入闭环 ─────────────────────────────────────────────────────

def test_import_all_four_formats():
    """4 个文件（含多格式）现按「每年完整报告」合并，不再因超 3 个报错。"""
    files = [
        ("files", ("利润表.xlsx", io.BytesIO(_xlsx_bytes()), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")),
        ("files", ("财报.docx", io.BytesIO(_build_docx_bytes()), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")),
        ("files", ("财报.pptx", io.BytesIO(_build_pptx_bytes()), "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
        ("files", ("利润表.csv", io.BytesIO("项目,2021,2022,2023\n营业收入,1,2,3\n".encode()), "text/csv")),
    ]
    r = CLIENT.post("/api/import", files=files)
    # 多文件合并成功；4 个文件各自解析后合并（含 csv 覆盖前 3 个文件的营业收入）
    assert r.status_code == 200, r.text
    body = r.json()
    assert len(body["previews"]) == 4


def test_import_xlsx_still_works():
    r = CLIENT.post("/api/import", files={"files": ("sample.xlsx", io.BytesIO(_xlsx_bytes()), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")})
    assert r.status_code == 200, r.text
    body = r.json()
    assert body["summary"]["years"] == [2021, 2022, 2023]
    assert body["previews"][0]["kind"] == "xlsx"


# ── 预览端点 ───────────────────────────────────────────────────────────

def test_preview_docx_content():
    r = CLIENT.post("/api/preview", files={"files": ("财报.docx", io.BytesIO(_build_docx_bytes()), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")})
    assert r.status_code == 200, r.text
    body = r.json()
    assert len(body["files"]) == 1
    f = body["files"][0]
    assert f["name"] == "财报.docx"
    assert f["kind"] == "docx"
    assert len(f["sections"]) == 3
    # 第一张表首行为项目+年份表头
    assert f["sections"][0]["grid"][0] == ["项目", "2021", "2022", "2023"]
    # 利润表科目行
    row_lookup = {row[0]: row for row in f["sections"][0]["grid"]}
    assert row_lookup["营业收入"][3] == "15200000"


def test_preview_pptx_content():
    r = CLIENT.post("/api/preview", files={"files": ("财报.pptx", io.BytesIO(_build_pptx_bytes()), "application/vnd.openxmlformats-officedocument.presentationml.presentation")})
    assert r.status_code == 200, r.text
    f = r.json()["files"][0]
    assert f["kind"] == "pptx"
    assert f["sections"][0]["grid"][0] == ["项目", "2021", "2022", "2023"]


def test_preview_csv_content():
    csv_text = "项目,2021,2022,2023\n营业收入,12000000,13500000,15200000\n"
    r = CLIENT.post("/api/preview", files={"files": ("利润表.csv", io.BytesIO(csv_text.encode()), "text/csv")})
    assert r.status_code == 200
    f = r.json()["files"][0]
    assert f["kind"] == "csv"
    assert f["sections"][0]["grid"][0] == ["项目", "2021", "2022", "2023"]


def test_preview_rejects_unsupported():
    r = CLIENT.post("/api/preview", files={"files": ("report.txt", io.BytesIO(b"hi"), "text/plain")})
    assert r.status_code == 400
    assert "仅支持" in r.json()["detail"]


def test_preview_bad_file_degrades():
    """损坏的 docx 预览应降级返回错误说明而非 500。"""
    r = CLIENT.post("/api/preview", files={"files": ("broken.docx", io.BytesIO(b"not a real docx"), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")})
    assert r.status_code == 200
    f = r.json()["files"][0]
    assert f["kind"] == "error"
    assert f["notes"][0].startswith("预览失败")


# ── PDF 解析与预览 ─────────────────────────────────────────────────────

PDF_MIME = "application/pdf"


def test_import_pdf():
    r = CLIENT.post("/api/import", files={"files": ("财报.pdf", io.BytesIO(_build_pdf_bytes()), PDF_MIME)}, data={"company_name": "PDF厂"})
    assert r.status_code == 200, r.text
    body = r.json()
    assert body["summary"]["company_name"] == "PDF厂"
    assert body["summary"]["years"] == [2021, 2022, 2023]
    data = session_mod.get_data()
    assert data is not None
    assert data.income_statement.get("营业收入", {}).get(2023) == 15_200_000
    assert body["previews"][0]["kind"] == "pdf"


def test_pdf_indicators_match_core():
    """PDF 指标与 core.finance 口径一致。"""
    r = CLIENT.post("/api/import", files={"files": ("财报.pdf", io.BytesIO(_build_pdf_bytes()), PDF_MIME)})
    assert r.status_code == 200, r.text
    ind2023 = r.json()["indicators"][2]
    assert ind2023["毛利率"]["value"] == 21.51
    assert ind2023["净利率"]["value"] == 2.43
    assert ind2023["增值税税负率"]["value"] == 1.92
    assert ind2023["增值税税负率"]["estimate"] is True


def test_preview_pdf_content():
    r = CLIENT.post("/api/preview", files={"files": ("财报.pdf", io.BytesIO(_build_pdf_bytes()), PDF_MIME)})
    assert r.status_code == 200, r.text
    f = r.json()["files"][0]
    assert f["kind"] == "pdf"
    assert f["name"] == "财报.pdf"
    # PDF 预览改为文本型：不再携带 base64 大图，按页采集文字
    assert f["images"] == [], "PDF 预览不再渲染图片"
    assert "pdf_type" in f, "PDF 预览应标注 pdf_type"
    assert any("第 1 页" in n for n in f["notes"])
    assert any("营业收入" in n for n in f["notes"])


def test_preview_pdf_is_text_not_image():
    """PDF 预览返回文字而非 base64 大图，杜绝超大响应拖垮浏览器。"""
    r = CLIENT.post("/api/preview", files={"files": ("财报.pdf", io.BytesIO(_build_pdf_bytes()), PDF_MIME)})
    f = r.json()["files"][0]
    assert not f["sections"], "PDF 不再用表格预览"
    assert f["images"] == [], "PDF 不再渲染图片"
    assert f["notes"], "PDF 预览应含文本/OCR 采样"
    joined = "\n".join(f["notes"])
    assert "data:image" not in joined
    assert "base64" not in joined


def test_import_pdf_scanned_degrades():
    """无文本层的扫描件 PDF：导入应返回 400 中文提示（引导配置 AI）而非 500。"""
    # 空 PDF 内容（无文本无表格）
    empty_pdf = b"%PDF-1.4\n1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]>>endobj\nxref\n0 4\n0000000000 65535 f \n0000000009 00000 n \n0000000052 00000 n \n0000000101 00000 n \ntrailer<</Size 4/Root 1 0 R>>\nstartxref\n171\n%%EOF"
    r = CLIENT.post("/api/import", files={"files": ("扫描件.pdf", io.BytesIO(empty_pdf), PDF_MIME)})
    assert r.status_code == 400
    detail = r.json()["detail"]
    assert "PDF" in detail
    assert "配置 AI" in detail
